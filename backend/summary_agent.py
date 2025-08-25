from fastapi import FastAPI, HTTPException
from pydantic import BaseModel
import os
from pinecone import Pinecone
from openai import OpenAI
import logging
from fastapi import FastAPI, File, UploadFile, HTTPException
from PyPDF2 import PdfReader
import os
from serpapi.google_search import GoogleSearch
from dotenv import load_dotenv
import zipfile
import tempfile
import pandas as pd
import docx
from typing import Union
import easyocr
from easyocr import Reader
import openai

# 加载环境变量
OPENAI_API_KEY = os.getenv("OPENAI_API_KEY")
PINECONE_API_KEY = os.getenv("PINECONE_API_KEY")
PINECONE_INDEX_NAME = os.getenv("PINECONE_INDEX_NAME")
SERPAPI_KEY = os.getenv("SERPAPI_KEY")

# 初始化 Pinecone 和 OpenAI
pc = Pinecone(api_key=PINECONE_API_KEY)
pinecone_index = pc.Index(PINECONE_INDEX_NAME)
openai_client = OpenAI(api_key=OPENAI_API_KEY)

# 初始化日志
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# 支持的文件扩展名
SUPPORTED_EXTENSIONS = ['.txt', '.pdf', '.docx', '.csv', '.xlsx', '.png', '.jpg', '.jpeg', '.py', '.zip', '.pdb', '.pptx']

app = FastAPI()

class SummaryRequest(BaseModel):
    facts: str
    issues: str
    reasoning: str
    decision: str

def get_embedding(text: str):
    """
    使用 OpenAI 嵌入 API 获取文本的嵌入向量。
    """
    response = openai_client.embeddings.create(
        input=text,
        model="text-embedding-ada-002"
    )
    return response.data[0].embedding

def extract_text_from_file(file_path: str) -> Union[str, None]:
    _, ext = os.path.splitext(file_path)
    ext = ext.lower()

    try:
        if ext == '.txt':
            with open(file_path, 'r', encoding='utf-8') as f:
                return f.read()

        elif ext == '.pdf':
            text = ""
            with open(file_path, 'rb') as f:
                reader_pdf = PdfReader(f)
                for page in reader_pdf.pages:
                    page_text = page.extract_text()
                    if page_text:
                        text += page_text + "\n"
            return text

        elif ext == '.docx':
            doc = docx.Document(file_path)
            return "\n".join([para.text for para in doc.paragraphs])

        elif ext == '.csv':
            df = pd.read_csv(file_path)
            return df.to_csv(index=False)

        elif ext == '.xlsx':
            df = pd.read_excel(file_path, engine='openpyxl')
            return df.to_csv(index=False)

        elif ext in ['.png', '.jpg', '.jpeg']:
            reader = Reader(['en'])
            result = reader.readtext(file_path, detail=0)
            return ' '.join(result)

        elif ext == '.py':
            with open(file_path, 'r', encoding='utf-8') as f:
                return f.read()

        elif ext == '.zip':
            return extract_text_from_zip(file_path)

        elif ext == '.pdb':
            with open(file_path, 'r', encoding='utf-8') as f:
                return f.read()

        elif ext == '.pptx':
            from pptx import Presentation
            prs = Presentation(file_path)
            text_runs = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_runs.append(shape.text)
            return "\n".join(text_runs)

        else:
            return f"Unsupported file type: {ext}"

    except Exception as e:
        logger.error(f"Error processing file {file_path}: {e}")
        return f"Error processing file: {e}"

def extract_text_from_zip(file_path: str) -> str:
    text = ""
    with zipfile.ZipFile(file_path, 'r') as zip_ref:
        with tempfile.TemporaryDirectory() as tmpdir:
            zip_ref.extractall(tmpdir)
            for root, _, files in os.walk(tmpdir):
                for file in files:
                    file_path_inner = os.path.join(root, file)
                    _, ext_inner = os.path.splitext(file_path_inner)
                    ext_inner = ext_inner.lower()
                    if ext_inner in SUPPORTED_EXTENSIONS:
                        extracted_text = extract_text_from_file(file_path_inner)
                        text += f"Extracted from {file}:\n{extracted_text}\n"
                    else:
                        text += f"Skipped unsupported file: {file}\n"
    return text

def generate_summary_file(prompt):
    """
    根据给定的提示词生成摘要报告。
    """
    try:
        # 调用 OpenAI Chat API
        response = openai.chat.completions.create(  # 保持您指定的方式
            model="gpt-4o",
            messages=[
                {
                    "role": "system",
                    "content": "This is a legal case. Please generate a summary report to help the lawyer understand the full case and details."
                },
                {
                    "role": "user",
                    "content": prompt
                }
            ],
            max_tokens=1500,
            temperature=0.7,
        )

        # 提取生成的内容
        if hasattr(response, "choices") and response.choices:
            advice = response.choices[0].message.content
            return advice
        else:
            return f"Unexpected response format: {response}"

    except Exception as e:
        return f"Error communicating with OpenAI: {e}"
    


def generate_summary_imput(prompt):
    """
    根据给定的提示词生成摘要报告。
    """
    try:
        # 调用 OpenAI Chat API
        response = openai.chat.completions.create(
            model="gpt-4o",
            messages=[
                {
                    "role": "system",
                    "content": "This is a legal case. Please generate a summary report to help the lawyer understand the full case and details."
                },
                {
                    "role": "user",
                    "content": prompt
                }
            ],
            max_tokens=1500,
            temperature=0.5,
        )

         # 提取生成的内容
        if hasattr(response, "choices") and response.choices:
            advice = response.choices[0].message.content
            return advice
        else:
            return f"Unexpected response format: {response}"

    except Exception as e:
        return f"Error communicating with OpenAI: {e}"


@app.post("/generate_summary")
async def generate_summary_endpoint(request: SummaryRequest):
    """
    接收用户输入，通过 Pinecone 匹配相似案例并生成摘要。
    """
    try:
        # 用户输入合并为一个文本
        user_input = f"Facts: {request.facts}\nIssues: {request.issues}\nReasoning: {request.reasoning}\nDecision: {request.decision}"
        logger.info(f"User input: {user_input}")

        # 获取嵌入向量
        user_embedding = get_embedding(user_input)
        logger.info("Generated embedding for user input.")

        # 查询 Pinecone 以获取最相似的案例
        search_results = pinecone_index.query(vector=user_embedding, top_k=1, include_metadata=True)
        logger.info(f"Search results: {search_results}")

        if "matches" not in search_results or len(search_results["matches"]) == 0:
            logger.warning("No similar cases found in Pinecone.")
            raise HTTPException(status_code=404, detail="No similar cases found.")

        # 获取最佳匹配案例
        best_match = search_results["matches"][0]
        metadata = best_match.get("metadata", {})
        case_plain_text = metadata.get("text", "").strip()

        if not case_plain_text:
            logger.warning(f"Plain text is empty for case ID: {best_match.get('id')}")
            raise HTTPException(status_code=400, detail="Similar case content is empty.")

        # 构建生成摘要的提示词
        prompt = f"""
        You are tasked with generating a summary for the following legal case.
        Context from a similar case:
        {case_plain_text}

        Generate a detailed and structured summary combining both contexts.
        """


        # 调用外部的 `generate_summary_imput` 函数
        summary = generate_summary_imput(prompt)

        # 检查生成的摘要
        if "Error communicating with OpenAI" in summary or "Unexpected response format" in summary:
            logger.error(f"OpenAI error: {summary}")
            raise HTTPException(status_code=500, detail=f"Failed to generate summary: {summary}")

        logger.info("Generated summary successfully.")
        return {"summary": summary}

    except Exception as e:
        logger.exception("Error generating summary.")
        raise HTTPException(status_code=500, detail=f"Error generating summary: {str(e)}")
    

@app.post("/upload_file")
async def upload_file(file: UploadFile = File(...)):
    try:
        # 保存上传的文件到临时目录
        temp_dir = tempfile.mkdtemp()
        file_path = os.path.join(temp_dir, file.filename)
        with open(file_path, "wb") as f:
            f.write(await file.read())

        # 获取文件扩展名
        _, ext = os.path.splitext(file.filename)
        ext = ext.lower()

        logger.info(f"Processing uploaded file: {file.filename} with extension: {ext}")

        # 如果是 .zip 文件，调用专用函数
        if ext == '.zip':
            file_content = extract_text_from_zip(file_path)
        else:
            # 处理其他文件类型
            file_content = extract_text_from_file(file_path)

        logger.info(f"Extracted file content length: {len(file_content)}")

        # 检查文件内容是否为空或不支持
        if not file_content or file_content.strip() == "":
            raise HTTPException(status_code=400, detail="Uploaded file has no readable text or unsupported file type.")

        # 生成提示词
        prompt = f"""
        Summarize the following document into the following categories:
        1. Case Summary
        2. Background
        3. Key Facts
        4. Key Issues
        5. Court's Opinion and Decision
        6. Key Observations
        7. Conclusion

        Document:
        {file_content}
        """
        logger.info(f"Generated prompt for OpenAI: {prompt[:500]}...")  # 仅打印部分内容，避免日志过长

        # 调用 OpenAI 模块生成报告
        report = generate_summary_file(prompt)

        # 检查生成结果
        if "Error communicating with OpenAI" in report or "Unexpected response format" in report:
            logger.error(f"OpenAI error: {report}")
            raise HTTPException(status_code=500, detail=f"Failed to generate summary: {report}")

        logger.info("Successfully generated summary report.")
        return {"report": report}

    except Exception as e:
        logger.error(f"Unexpected error during file processing: {e}")
        raise HTTPException(status_code=500, detail=f"Unexpected error: {str(e)}")